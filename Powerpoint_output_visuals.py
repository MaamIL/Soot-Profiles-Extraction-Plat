import os
import re
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor  # Import RGBColor
from PIL import Image
from Logger import CustomLogger

class PowerPointVisual:
    def __init__(self, config, train_losses, val_losses, test_loss):
        self.logger = CustomLogger(config.log_filename, self.__class__.__name__).get_logger()
        self.folder = config.out_dir
        self.train_losses = train_losses 
        self.val_losses =  val_losses
        self.test_loss = test_loss

    def collect_images(self):
        self.logger.info(f"Collecting images from {self.folder}")
        try:
            pattern = re.compile(r'^(\d+)_(fv|T)_E_(\d+)_heatmap\.jpg$')
            images = defaultdict(lambda: {'input': None, 'fv': {}, 'T': {}, 'target_fv': None, 'target_T': None})
            
            for file in os.listdir(self.folder):
                # match = re.match(r'losses\.png$', file)
                # if match:
                #     images[match.group(1)]['losses'] = file
                #     self.logger.info(f"Found losses image: {file}")
                #     continue

                match = re.match(r'^(\d+)_Input\.jpg$', file)
                if match:
                    images[match.group(1)]['input'] = file
                    self.logger.info(f"Found input image: {file}")
                    continue
                
                match = pattern.match(file)
                if match:
                    num, attr, time = match.groups()
                    images[num][attr][int(time)] = file
                    self.logger.info(f"Found {attr} image: {file}")
                    continue
                
                match = re.match(r'^(\d+)_(fv|T)_Target_heatmap\.jpg$', file)
                if match:
                    num, attr = match.groups()
                    images[num][f'target_{attr}'] = file
                    self.logger.info(f"Found target {attr} image: {file}")
                    
                
                    
        except Exception as e:
            self.logger.error(f"Error in collect_images: {e}")

        return images

    def image_proportions(self, img_path):    
        with Image.open(img_path) as img:
            width, height = img.size        
            ratio = height / width  
            return ratio
        
    def create_presentation(self, images, losses, output_file="EpochsPresentation.pptx"):
        prs = Presentation()

        # First slide with losses image and caption
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        title = slide.shapes.title
        title.text = "Losses"
        center = Inches(1.5)
        img_width = Inches(7)
        img_height = Inches(3)
        img = slide.shapes.add_picture(os.path.join(self.folder, losses), center, center, img_width, img_height)
        #add image under this one
        img_width = img_height / self.image_proportions(os.path.join(self.folder, losses))
        img = slide.shapes.add_picture(os.path.join(self.folder, "losses_0_20.png"), center, center + img_height + Inches(0.5), img_width, img_height)
        caption = slide.shapes.add_textbox(center, center + img_height + Inches(0.1), img_width, Inches(0.5))
        text_frame = caption.text_frame
        text_frame.text = f"Train Loss: {self.train_losses[-1]} Val Loss: {self.val_losses[-1]} Test Loss: {self.test_loss}"
        text_frame.paragraphs[0].font.size = Inches(0.3)

        #other slides with images of epochs changes
        for num, data in sorted(images.items(), key=lambda x: int(x[0])):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Sample #{num}"
            
            left = Inches(0.5)
            top = Inches(1.3)
            center = Inches(1.5)
            img_width = Inches(2.5)
            img_height = Inches(2.5)

            if data['input']:
                img_width = img_height*2 / self.image_proportions(os.path.join(self.folder, data['input']))
                slide.shapes.add_picture(os.path.join(self.folder, data['input']), left, center, img_width, img_height*2)
            
            start_x = left + img_width + Inches(0.25)
            y_offset = [top, top + img_height + Inches(0.5)]
            
            times = [0, 49, 99, 149]
            for i, attr in enumerate(['fv', 'T']):
                    x = start_x
                    for t in times:
                        if t in data[attr]:
                            img_width = img_height / self.image_proportions(os.path.join(self.folder, data[attr][t]))
                            img = slide.shapes.add_picture(os.path.join(self.folder, data[attr][t]), x, y_offset[i], img_width, img_height)
                            caption = slide.shapes.add_textbox(x, y_offset[i] + img_height + Inches(0.1), img_width, Inches(0.5))
                            text_frame = caption.text_frame
                            text_frame.text = f"{attr} epoch {t}"
                            text_frame.paragraphs[0].font.size = Inches(0.3)
                        x += img_width + Inches(0.1)
                    
                    # Add vertical line before target
                    line = slide.shapes.add_shape(1, x - Inches(0.1), y_offset[i], Inches(0.05), img_height)
                    line.line.color.rgb = RGBColor(0, 0, 0)  # Black color
                    
                    if data[f'target_{attr}']:
                        img_width = img_height / self.image_proportions(os.path.join(self.folder, data[f'target_{attr}']))
                        img = slide.shapes.add_picture(os.path.join(self.folder, data[f'target_{attr}']), x, y_offset[i], img_width, img_height)
                        caption = slide.shapes.add_textbox(x, y_offset[i] + img_height + Inches(0.1), img_width, Inches(0.3))
                        text_frame = caption.text_frame
                        text_frame.text = f"{attr} target"
                        text_frame.paragraphs[0].font.size = Inches(0.3)

        prs.save(os.path.join(self.folder, output_file))

      
    # folder = "Heatmaps_2025-0313-184458"
    # images = collect_images(folder)
    # create_presentation(images, folder)
